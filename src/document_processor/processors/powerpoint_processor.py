#!filepath: powerpoint_processor.py
"""
PowerPoint (PPT/PPTX) processor for RAG Generator Tool.
Converts PowerPoint presentations to Markdown format with images.
"""

import os
import hashlib
import logging
from typing import List, Tuple, Dict, Any, Optional
from ..config import ConverterConfig
from .base_processor import BaseDocumentProcessor

logger = logging.getLogger(__name__)

class PowerPointProcessor(BaseDocumentProcessor):
    """Processor for PowerPoint (PPT/PPTX) files"""
    
    def __init__(self, config: ConverterConfig) -> None:
        super().__init__(config)
    
    @classmethod
    def get_supported_extensions(cls) -> List[str]:
        """Return the file extensions supported by this processor"""
        return [".pptx", ".ppt"]
    
    @classmethod
    def can_process(cls, file_path: str) -> bool:
        """PowerPoint specific filtering"""
        # Base class already handles ~$ files
        return super(PowerPointProcessor, cls).can_process(file_path)
    
    def process(self, file_path: str, output_dir: str, media_dir: str) -> Tuple[Optional[str], Dict[str, Any]]:
        """
        Process PowerPoint file and convert to Markdown
        
        Args:
            file_path: Path to the PowerPoint file
            output_dir: Directory for temporary extraction
            media_dir: Directory to save extracted media
            
        Returns:
            tuple: (markdown_content, metadata_dict)
        """
        # Handle old PPT format with message
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".ppt":
            logger.info(f"Processing older PPT format file: {os.path.basename(file_path)}")
            
            # For PPT files, add a notice to the markdown
            ppt_notice = (
                "# PPT Format Notice\n\n"
                "This file is in the older PPT format, which has limited support in this tool. "
                "For best results, please consider converting it to PPTX format using Microsoft PowerPoint "
                "or another compatible application.\n\n"
                "Some content or formatting may not be properly extracted from this file.\n\n"
                "---\n\n"
            )
              # We'll still try to process it as best we can
            md_content, metadata = self.process_powerpoint_file(file_path, media_dir)
            
            # Add the notice at the beginning
            if md_content and not md_content.startswith("# Error"):
                md_content = ppt_notice + md_content
                
            return md_content, metadata
        
        # Process PPTX files normally
        return self.process_powerpoint_file(file_path, media_dir)
    def process_powerpoint_file(self, file_path: str, media_dir: str) -> Tuple[Optional[str], Dict[str, Any]]:
        """
        Process PowerPoint (PPTX) files and convert to Markdown with images.
        
        Args:
            file_path: Path to the PowerPoint file
            media_dir: Directory to save extracted media
            
        Returns:
            tuple: (markdown_content, metadata_dict)
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            logger.error("python-pptx library not installed. Install with: pip install python-pptx")
            return None, {
                'source_filename': os.path.basename(file_path),
                'parser': 'pptx_parser',
                'error': 'Missing dependency: python-pptx'
            }
        logger.info(f"Processing PowerPoint file: {file_path}")
        
        # Create media directory if it doesn't exist
        os.makedirs(media_dir, exist_ok=True)
        
        try:
            # Open the presentation
            presentation = Presentation(file_path)
            
            # Initialize markdown content
            md_content = []
            image_count = 0
            
            # Get presentation metadata
            metadata = {
                'source_filename': os.path.basename(file_path),
                'parser': 'pptx_parser',
                'slide_count': len(presentation.slides),
                'title': self.get_presentation_title(presentation)
            }
            
            # Add title to markdown
            if metadata['title']:
                md_content.append(f"# {metadata['title']}\n")
            
            # Process each slide
            for slide_index, slide in enumerate(presentation.slides, 1):
                
                # Add slide heading
                slide_title = self.get_slide_title(slide) or f"Slide {slide_index}"
                md_content.append(f"## {slide_title}\n")
                
                # Process text content (shapes, text boxes, etc.)
                text_content = self.extract_slide_text(slide)
                if text_content:
                    md_content.append(text_content)
                
                # Process images in the slide
                for shape in slide.shapes:
                    try:
                        # Only Picture shapes have image data in python-pptx
                        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                            img = getattr(shape, "image", None)
                            if not img:
                                continue
                            # Extract image data
                            image_bytes = img.blob
                            
                            # Generate a unique filename
                            image_hash = hashlib.md5(image_bytes).hexdigest()[:12]
                            image_ext = self.get_image_extension(image_bytes) or '.png'
                            image_filename = f"ppt_img_{image_count}_{image_hash}{image_ext}"
                            image_path = os.path.join(media_dir, image_filename)
                            
                            # Save the image
                            with open(image_path, 'wb') as img_file:
                                img_file.write(image_bytes)
                            
                            # Add image reference to markdown
                            alt_text = (
                                getattr(shape, 'alt_text', None)
                                or getattr(shape, 'name', None)
                                or f"Slide {slide_index} Image {image_count}"
                            )
                            md_content.append(f"\n![{alt_text}](media/{image_filename})\n")
                            
                            image_count += 1
                    except Exception as e:
                        logger.warning(f"Failed to process image in slide {slide_index}: {e}")
                
                # Add slide notes if any
                notes_text = self.extract_slide_notes(slide)
                if notes_text:
                    md_content.append(f"\n> **Slide Notes:** {notes_text}\n")
                
                # Add separator between slides
                md_content.append("\n---\n")
            
            # Join all content into a single string
            final_markdown = "\n".join(md_content)
            
            # Update metadata
            metadata['image_count'] = image_count
            
            return final_markdown, metadata
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file {file_path}: {e}", exc_info=True)
            return None, {
                'source_filename': os.path.basename(file_path),
                'parser': 'pptx_parser',
                'error': str(e)
            }

    def get_presentation_title(self, presentation: Any) -> Optional[str]:
        """Extract title from presentation properties"""
        if hasattr(presentation.core_properties, 'title') and presentation.core_properties.title:
            return presentation.core_properties.title
        
        # Try to get title from first slide
        if presentation.slides and len(presentation.slides) > 0:
            return self.get_slide_title(presentation.slides[0])
        
        return None

    def get_slide_title(self, slide: Any) -> Optional[str]:
        """Extract title from a slide"""
        for shape in slide.shapes:
            # Safe check for title shapes - approach 1: text and is_title property
            if hasattr(shape, 'text') and hasattr(shape, 'is_title') and shape.is_title:
                return shape.text.strip()
            
            # Safe check for title shapes - approach 2: shape name contains "title"
            if hasattr(shape, 'name') and 'title' in shape.name.lower() and hasattr(shape, 'text'):
                return shape.text.strip()
                
            # Safe check for placeholder type - this is the problematic part
            try:
                # Use a try-except block instead of hasattr for placeholder_format
                if hasattr(shape, 'placeholder_format'):
                    placeholder_type = shape.placeholder_format.type
                    if placeholder_type == 1 and hasattr(shape, 'text'):  # 1 = title placeholder
                        return shape.text.strip()
            except ValueError:
                # Skip shapes that raise ValueError when accessing placeholder_format
                pass
        
        return None

    def extract_slide_text(self, slide: Any) -> str:
        """Extract formatted text content from a slide"""
        text_content = []
        
        for shape in slide.shapes:
            # First check if it has text and the text is not empty
            if hasattr(shape, 'text') and shape.text.strip():
                # Then safely check if this is a title (only if the attribute exists)
                is_title = hasattr(shape, 'is_title') and shape.is_title
                
                # Only process non-title text
                if not is_title:
                    # Process text based on shape type and level
                    text = shape.text.strip()
                    
                    # Handle bullet points if available
                    if hasattr(shape, 'paragraphs'):
                        for i, paragraph in enumerate(shape.paragraphs):
                            if hasattr(paragraph, 'level'):
                                # Indent based on level
                                indent = '  ' * paragraph.level
                                bullet = '* ' if paragraph.level == 0 else '- '
                                text_content.append(f"{indent}{bullet}{paragraph.text.strip()}")
                            else:
                                text_content.append(paragraph.text.strip())
                    else:
                        text_content.append(text)
        
        return "\n".join(text_content) if text_content else ""

    def extract_slide_notes(self, slide: Any) -> Optional[str]:
        """Extract notes from a slide"""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            return slide.notes_slide.notes_text_frame.text.strip()
        return None

    def get_image_extension(self, image_bytes: bytes) -> Optional[str]:
        """Determine file extension based on image data"""
        # Try using Pillow for robust format detection (compatible with Python 3.13+)
        try:
            from io import BytesIO
            from PIL import Image
            with Image.open(BytesIO(image_bytes)) as img:
                fmt = (img.format or '').lower()
                if fmt == 'jpeg':
                    return '.jpg'
                if fmt:
                    return f'.{fmt}'
        except Exception:
            # Fall back to magic byte checks below
            pass
    
        # Fallback: Check magic bytes manually for common formats
        if image_bytes.startswith(b'\xFF\xD8'):
            return '.jpg'
        elif image_bytes.startswith(b'\x89PNG\r\n\x1A\n'):
            return '.png'
        elif image_bytes.startswith(b'GIF87a') or image_bytes.startswith(b'GIF89a'):
            return '.gif'
        elif image_bytes.startswith(b'\x42\x4D'):
            return '.bmp'
        elif image_bytes.startswith((b'II*\x00', b'MM\x00*')):
            return '.tif'
        elif image_bytes.startswith(b'RIFF') and image_bytes[8:12] == b'WEBP':
            return '.webp'
    
        # Default to PNG if can't determine
        return '.png'
